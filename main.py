import os
import sys
import csv
import fitz
from loguru import logger
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import openpyxl

PASTA_CONTAINER = "/converter"
PASTA_TXT = "/converter/txt"
FORMATOS_SUPORTADOS = {".pdf", ".docx", ".xlsx", ".csv", ".pptx", ".html", ".htm"}

logger.remove()
logger.add(
    sys.stdout,
    format="<green>{time:YYYY-MM-DD HH:mm:ss}</green> | <level>{level: <8}</level> | {message}",
    level="DEBUG",
    colorize=True,
)


def buscar_arquivos(pasta):
    logger.info(f"Varrendo pasta: {pasta}")

    if not os.path.isdir(pasta):
        logger.error(f"Pasta não encontrada: {pasta}")
        sys.exit(1)

    arquivos = [
        os.path.join(pasta, f)
        for f in os.listdir(pasta)
        if os.path.splitext(f)[1].lower() in FORMATOS_SUPORTADOS
        and os.path.isfile(os.path.join(pasta, f))
    ]

    if not arquivos:
        logger.warning(f"Nenhum arquivo suportado encontrado em: {pasta}")
        sys.exit(0)

    logger.info(f"{len(arquivos)} arquivo(s) encontrado(s).")
    return arquivos


def converter_pdf(caminho_origem, arquivo_final):
    try:
        pdf = fitz.open(caminho_origem)
    except Exception as e:
        logger.error(f"Falha ao abrir PDF: {e}")
        return

    total_paginas = len(pdf)
    logger.info(f"Total de páginas: {total_paginas}")

    if total_paginas == 0:
        logger.warning("PDF sem páginas — pulando.")
        pdf.close()
        return

    paginas_sem_texto = 0

    try:
        with open(arquivo_final, "w", encoding="utf-8") as txt:
            for i in range(total_paginas):
                texto = pdf[i].get_text()
                if not texto.strip():
                    paginas_sem_texto += 1
                    logger.warning(f"Página {i + 1}/{total_paginas} sem texto detectável.")
                else:
                    logger.debug(f"Página {i + 1}/{total_paginas} extraída ({len(texto)} caracteres).")
                txt.write(texto)
                txt.write("\n" + "-" * 20 + "\n")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")
        return
    finally:
        pdf.close()

    if paginas_sem_texto > 0:
        logger.warning(
            f"{paginas_sem_texto}/{total_paginas} página(s) sem texto — "
            "considere usar OCR para documentos escaneados."
        )


def converter_docx(caminho_origem, arquivo_final):
    try:
        doc = Document(caminho_origem)
    except Exception as e:
        logger.error(f"Falha ao abrir DOCX: {e}")
        return

    try:
        with open(arquivo_final, "w", encoding="utf-8") as txt:
            for para in doc.paragraphs:
                if para.text.strip():
                    txt.write(para.text + "\n")

            for tabela in doc.tables:
                txt.write("\n" + "-" * 20 + "\n")
                for linha in tabela.rows:
                    conteudo = "\t".join(cell.text.strip() for cell in linha.cells)
                    txt.write(conteudo + "\n")

        logger.debug(f"{len(doc.paragraphs)} parágrafos extraídos.")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")


def converter_excel(caminho_origem, arquivo_final):
    try:
        wb = openpyxl.load_workbook(caminho_origem, data_only=True)
    except Exception as e:
        logger.error(f"Falha ao abrir Excel: {e}")
        return

    try:
        with open(arquivo_final, "w", encoding="utf-8") as txt:
            for nome_aba in wb.sheetnames:
                ws = wb[nome_aba]
                logger.debug(f"Extraindo aba: {nome_aba}")
                txt.write(f"=== Aba: {nome_aba} ===\n")
                for linha in ws.iter_rows(values_only=True):
                    conteudo = "\t".join(str(cell) if cell is not None else "" for cell in linha)
                    if conteudo.strip():
                        txt.write(conteudo + "\n")
                txt.write("\n" + "-" * 20 + "\n")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")


def converter_csv(caminho_origem, arquivo_final):
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            with open(caminho_origem, "r", encoding=encoding, newline="") as f:
                reader = csv.reader(f)
                linhas = list(reader)
            break
        except (UnicodeDecodeError, Exception):
            continue
    else:
        logger.error(f"Não foi possível detectar encoding do CSV: {caminho_origem}")
        return

    try:
        with open(arquivo_final, "w", encoding="utf-8") as txt:
            for linha in linhas:
                txt.write("\t".join(linha) + "\n")
        logger.debug(f"{len(linhas)} linhas extraídas do CSV.")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")


def converter_pptx(caminho_origem, arquivo_final):
    try:
        prs = Presentation(caminho_origem)
    except Exception as e:
        logger.error(f"Falha ao abrir PPTX: {e}")
        return

    try:
        with open(arquivo_final, "w", encoding="utf-8") as txt:
            for i, slide in enumerate(prs.slides):
                txt.write(f"=== Slide {i + 1} ===\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        txt.write(shape.text + "\n")
                txt.write("\n" + "-" * 20 + "\n")
        logger.debug(f"{len(prs.slides)} slides extraídos.")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")


def converter_html(caminho_origem, arquivo_final):
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            with open(caminho_origem, "r", encoding=encoding) as f:
                conteudo = f.read()
            break
        except UnicodeDecodeError:
            continue
    else:
        logger.error(f"Não foi possível ler o HTML: {caminho_origem}")
        return

    try:
        soup = BeautifulSoup(conteudo, "lxml")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        texto = soup.get_text(separator="\n")
        linhas = [l.strip() for l in texto.splitlines() if l.strip()]

        with open(arquivo_final, "w", encoding="utf-8") as txt:
            txt.write("\n".join(linhas))
        logger.debug(f"{len(linhas)} linhas extraídas do HTML.")
    except IOError as e:
        logger.error(f"Falha ao escrever '{arquivo_final}': {e}")


DISPATCH = {
    ".pdf": converter_pdf,
    ".docx": converter_docx,
    ".xlsx": converter_excel,
    ".csv": converter_csv,
    ".pptx": converter_pptx,
    ".html": converter_html,
    ".htm": converter_html,
}


def converter(caminho_origem):
    ext = os.path.splitext(caminho_origem)[1].lower()
    nome_base = os.path.splitext(os.path.basename(caminho_origem))[0]
    arquivo_final = os.path.join(PASTA_TXT, nome_base + ".txt")

    logger.info(f"Convertendo: {os.path.basename(caminho_origem)}")

    func = DISPATCH.get(ext)
    if not func:
        logger.warning(f"Formato não suportado: {ext} — pulando.")
        return

    func(caminho_origem, arquivo_final)
    logger.success(f"Salvo em: {arquivo_final}")


def main():
    logger.info("Iniciando conversor → TXT")
    logger.debug(f"Pasta alvo: {PASTA_CONTAINER}")
    logger.debug(f"Pasta de saída: {PASTA_TXT}")
    logger.debug(f"Formatos suportados: {', '.join(sorted(FORMATOS_SUPORTADOS))}")

    os.makedirs(PASTA_TXT, exist_ok=True)

    arquivos = buscar_arquivos(PASTA_CONTAINER)

    for arquivo in arquivos:
        converter(arquivo)

    logger.info("Processo finalizado.")


if __name__ == "__main__":
    main()
